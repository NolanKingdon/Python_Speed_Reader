import os 
import time
from time import gmtime, strftime
import glob
from pptx import Presentation
from docx import Document
import PyPDF2 as pypdf2

'''

MODEL - Sends information to the view

'''

class Reader:
    def __init__(self, WPM, clearKeyword):
        self.WPM = WPM
        self._READ_SPEED = (60/WPM)/2 #The pauses are cumulative, and add lag. Adjust for it here
        self.__CLEAR = clearKeyword
        self._TOTAL_WORDS = 0
        self._TIME_STARTED = ""
        self._LAST_WORD = ""
        self._CURRENT_PARAGRAPH = ""
        self._CURRENT_FILE = ""
        
    # Opens the folder provided by filepath. Returns the fileList OR False if address is bad
    def openFolder(self, filePath):
        fileList = []
        try:
            for file in os.listdir(filePath):
                if filePath[-1] == "\\":
                    fileList.append(filePath + file)
                else: 
                    fileList.append(filePath + "\\" + file)
            return fileList
        except:
            #Signalling an error
            print(err)
            return False
        
    # Loads in sources from file locations
    def loadSources(self, fileLocations):
        #os.system(self.__CLEAR)
        sources = {}
        for file in fileLocations:
            ext = file.split(".")[-1].lower()
            if (ext == "pptx"):
                sources[file] = self.__scrapePPTX(file)
            if (ext == "docx"):
                sources[file] = self.__scrapeDOCX(file)
            if (ext == "txt"):
                sources[file] = self.__scrapeTXT(file)
            if (ext == "pdf"):
                sources[file] = self.__scrapePDF(file)
        return sources

    # Counts in the application - Will move to view in future
    def setStartTime(self):
        self._TIME_STARTED = strftime("%H:%M:%S", gmtime()) # After we count in, set the start time

    # Gets the file path and the total words in the individual source file
    def parseFiles(self, file, sources):
        splitPath = file.split("/") if file[-1] == "/" else file.split("\\")
        totalWords = self.__getWordCount(sources[file])
        return (splitPath[-1], totalWords)

    def setReadSpeed(self, WPM):
        self._READ_SPEED = (60/WPM)/2

    # Getter for last read word
    def getLastWord(self):
        return self._LAST_WORD

    # Clearing our last word for when we're done the read or reset the reader
    def clearLastWord(self):
        self._LAST_WORD = "" 
    
    # Getter for read speed - value we time.sleep by
    def getReadSpeed(self):
        return self._READ_SPEED            

    # Getter for total words of all sources
    def getTotalWords(self):
        return self._TOTAL_WORDS

    # Getter for time started
    def getStartTime(self):
        return self._TIME_STARTED

    # Getter to calculate time since start
    def getElapsedTime(self, startTime):
        #Start time will inherently be 0:0:00 and we will increment using maths and logic
        startList = startTime.split(":")
        currentTime = strftime("%H:%M:%S", gmtime())
        currentTime = currentTime.split(":")
        hoursPassed = str(int(currentTime[0]) - int(startList[0]))
        minutesPassed = str(int(currentTime[1]) - int(startList[1]))
        secondsPassed = str(int(currentTime[2]) - int(startList[2]) % (int(currentTime[2])+ 1))
        elapsedTime = hoursPassed + ":" + minutesPassed + ":" + secondsPassed
        return elapsedTime

    # Getter for actual calculated WPM score (Accounts for punctuation and slows)
    def getTrueWPM(self, elapsedTime):
        calculatedWPM = self.__calculateWPM(elapsedTime, self._TOTAL_WORDS)
        return calculatedWPM

    def __scrapePDF(self, filePath):
        paragraphs = []
        pdfReader = pypdf2.PdfFileReader(filePath)
        pageNo = pdfReader.getNumPages()
        for i in range(0, pageNo):
            paragraphs += pdfReader.getPage(i).extractText().split("\n")
        return paragraphs

    def __scrapeTXT(self, filePath):
        paragraphs = []
        file = open(filePath, 'r')
        unformattedText = file.read();
        paragraphs = unformattedText.split("\n") #Splitting by paragraph
        return paragraphs

    def __scrapeDOCX(self, filePath):
        paragraphs = []
        doc = Document(filePath)
        for paragraph in doc.paragraphs:
            paragraphs.append(paragraph.text)
        return paragraphs

    def __scrapePPTX(self, filePath):
        paragraphs = []
        pres = Presentation(filePath)
        for slide in pres.slides: #Looking through our slides
            for shape in slide.shapes: #Looking through objects in that slide
                if not shape.has_text_frame: #Seeing if we have any text on that slide
                    continue
                #for title in shape.name
                for paragraph in shape.text_frame.paragraphs: #reading all paragraphs on the slide
                    paragraphs.append(paragraph.text)
        return paragraphs

    def __getWordCount(self, paragraphs):
        wordCount = 0
        for paragraph in paragraphs:
            for word in paragraph.split(" "):
                wordCount = wordCount + 1
        return wordCount

    def __calculateWPM(self, elapsedTime, totalWords):
        splitTime = elapsedTime.split(":")
        totalSeconds = (int(splitTime[0]) * 60 * 60) + (int(splitTime[1]) * 60) + int(splitTime[2])    
        return round((totalWords/totalSeconds)*60)

    # Determines center word and sleep modifier of specific word
    def read(self, source, word):
        #Finding the letter roughly one third of the way through the word to center
        centerPos = 1
        origin = source.split("/") if source[-1] == "/" else source.split("\\")
        #Setting sleep Modifier based on word length or placing in sentence
        sleepModifier = 1 #How long we want to wait on each word
        if (len(word) > 0): #Ensuring we don't run into issues with blank spaces       
            if(len(word) > 10):
                sleepModifier = 1.5
            elif(word[-1] in ".,-?!;:"):
                sleepModifier = 1.2
            centerPos = int(len(word)/3)
        #Tracking our last word
        self._LAST_WORD = word            
        # Incrementing our total words
        self._TOTAL_WORDS = self._TOTAL_WORDS + 1
        # Tuples are great.
        return (centerPos, sleepModifier)
                
