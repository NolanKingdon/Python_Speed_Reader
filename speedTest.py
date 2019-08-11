import os 
import time
from time import gmtime, strftime
import glob
from pptx import Presentation
from docx import Document
import PyPDF2 as pypdf2

# NOTE: This is just a quick and dirty tool I made because I needed something to help study.
#       In the future I am going to be redoing this to be better form
#       That means having an actual file structure and MVC approach

def openFolder():
    fileList = []
    filePath = input("Enter folder you want to read\n")
    for file in os.listdir(filePath):
        if filePath[-1] == "\\":
            fileList.append(filePath + file)
        else: 
            fileList.append(filePath + "\\" + file)
    return fileList

def scrapePDF(filePath):
    paragraphs = []
    pdfReader = pypdf2.PdfFileReader(filePath)
    pageNo = pdfReader.getNumPages()
#    print(pdfReader.getPage(0).extractText())

    for i in range(0, pageNo):
        paragraphs += pdfReader.getPage(i).extractText().split("\n")

    return paragraphs


def scrapeTXT(filePath):
    paragraphs = []
    file = open(filePath, 'r')
    unformattedText = file.read();
    paragraphs = unformattedText.split("\n") #Splitting by paragraph
    return paragraphs

def scrapeDOCX(filePath):
    paragraphs = []
    doc = Document(filePath)
    for paragraph in doc.paragraphs:
        paragraphs.append(paragraph.text)
    return paragraphs

def scrapePPTX(filePath):
    paragraphs = []
    pres = Presentation(filePath)
    for slide in pres.slides: #Looking through our slides
        for shape in slide.shapes: #Looking through objects in that slide
            if not shape.has_text_frame: #Seeing if we have any text on that slide
                continue
            #for title in shape.name
            for paragraph in shape.text_frame.paragraphs: #reading all paragraphs on the slide
                paragraphs.append(paragraph.text)
    return paragraphs

def getWordCount(paragraphs):
    wordCount = 0
    for paragraph in paragraphs:
        for word in paragraph.split(" "):
            wordCount = wordCount + 1
    return wordCount

def getElapsedTime(startTime):
    #Start time will inherently be 0:0:00 and we will increment using maths and logic
    startList = startTime.split(":")
    currentTime = strftime("%H:%M:%S", gmtime())
    currentTime = currentTime.split(":")
    hoursPassed = str(int(currentTime[0]) - int(startList[0]))
    minutesPassed = str(int(currentTime[1]) - int(startList[1]))
    secondsPassed = str(int(currentTime[2]) - int(startList[2]) % (int(currentTime[2])+ 1))
    
    elapsedTime = hoursPassed + ":" + minutesPassed + ":" + secondsPassed
    return elapsedTime

def calculateWPM(elapsedTime, totalWords):
    splitTime = elapsedTime.split(":")
    totalSeconds = (int(splitTime[0]) * 60 * 60) + (int(splitTime[1]) * 60) + int(splitTime[2])    
    return round((totalWords/totalSeconds)*60)


def countIn():
    print("Starting program in 3")
    print("\n\n\n\n" + " " * 12 + RED + "X" + RESET + "<----- Look here") 
    time.sleep(1)
    os.system(CLEAR)
    print("Starting program in 2")
    print("\n\n\n\n" + " " * 12 + RED + "X" + RESET + "<----- Look here") 
    time.sleep(1)
    os.system(CLEAR)
    print("Starting program in 1")
    print("\n\n\n\n" + " " * 12 + RED + "X" + RESET + "<----- Look here")
    time.sleep(1)
    os.system(CLEAR)

def read(source, paragraphs, sourceWords):
    count = 0
    global TOTAL_WORDS # Change this later - bad practice
    for paragraph in paragraphs:
        for word in paragraph.split(" "):
            #Finding the letter roughly one third of the way through the word to center
            centerPos = int(len(word)/3)
            origin = source.split("/") if source[-1] == "/" else source.split("\\")
            #Setting sleep Modifier based on word length or placing in sentence
            sleepModifier = 1 #How long we want to wait on each word        
            if (len(word) > 0): #Ensuring we don't run into issues with blank spaces       
                if(len(word) > 7):
                    sleepModifier = 1.5
                elif(word[-1] in ".,-?!;:"):
                    sleepModifier = 1.2
                        
                formattedWord = word[0:centerPos] + RED + word[centerPos] + RESET + word[centerPos+1:len(word)] + RESET        
                print("Origin: " + origin[-1]) #Keep in mind you'll want to track slide/file origin here
                print("\n\n\n\n" + " "*(12 - centerPos) + formattedWord)
                print("\n\n\n\n")
                print(str(count) + " / " + str(sourceWords))
                time.sleep(READ_SPEED * sleepModifier) #Waiting changes depending on context
                count = count + 1
            else:
                print("Origin: " + origin[-1]) #Keep in mind you'll want to track slide/file origin here
                print("\n\n\n\n")
                print(str(count) + " / " + str(sourceWords))
                count = count + 1
            
            TOTAL_WORDS = TOTAL_WORDS + 1
            os.system(CLEAR) #clearing our terminal so we can stay at the top of the window
            

# VARIABLE DECLARATIONS

CLEAR = "cls" if os.name == "nt" else "clear"
READ_SPEED = 0.09 #0.07 is manageable
WPM = round(60 / READ_SPEED) #Seconds in a minute.
RED = "\033[1;31m" #Formatting options
RESET = "\033[0;0m" #Formatting options
BOLD = "\033[1m" #Formatting options
TOTAL_WORDS = 0
TIME_STARTED = "0"
TIME_END = "0"

os.system(CLEAR) # Getting all the junk out of the way before we start

#  --- MAIN RUNTIME ---

fileLocations = openFolder()
sources = {}

for file in fileLocations:
    ext = file.split(".")[-1].lower()
#    print(ext)
    print("Reading " + file + "...")
    time.sleep(0.5)
    os.system(CLEAR)
    if (ext == "pptx"):
        sources[file] = scrapePPTX(file)
    if (ext == "docx"):
        sources[file] = scrapeDOCX(file)
    if (ext == "txt"):
        sources[file] = scrapeTXT(file)
    if (ext == "pdf"):
        sources[file] = scrapePDF(file)
print("Starting ...")
time.sleep(1)
os.system(CLEAR)
countIn()

TIME_STARTED = strftime("%H:%M:%S", gmtime())
for file in fileLocations:
    splitPath = file.split("/") if file[-1] == "/" else file.split("\\")
    print("\n\n\n\n\n" + " "*12 + RED + "N" + RESET + "ext Source: " + splitPath[-1])
    time.sleep(2)
    os.system(CLEAR)
    totalWords = getWordCount(sources[file])
    read(file, sources[file], totalWords)
    os.system(CLEAR)

print("Done")
print("Time Started: " + TIME_STARTED)
print("Time Reading: " + getElapsedTime(TIME_STARTED))
print("Time Ended: " + strftime("%H:%M:%S", gmtime()))
print("Total Words Read: " + str(TOTAL_WORDS))
print("Read speed: " + str(calculateWPM(getElapsedTime(TIME_STARTED), TOTAL_WORDS)) + "WPM")


